#!/usr/bin/env python3
"""AI Buddy デモ紹介デッキ（日本語 / Glean ブランド）を .pptx で生成する。

生成物は編集可能な PowerPoint。Google Slides にもそのままインポートできる。
usage: .venv_slides/bin/python scripts/build_deck.py [out.pptx]
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "frontend" / "public" / "logos" / "glean-logo.png"

# ── Glean ブランド・パレット ────────────────────────────────────────────
BLUE = RGBColor(0x34, 0x3C, 0xED)      # Glean プライマリ
BLUE_DK = RGBColor(0x21, 0x25, 0xA8)
INK = RGBColor(0x0F, 0x10, 0x35)       # 見出し
BODY = RGBColor(0x3D, 0x40, 0x5C)      # 本文
MUTED = RGBColor(0x76, 0x79, 0x96)
TINT = RGBColor(0xEE, 0xEF, 0xFE)      # 薄紫の面
TINT2 = RGBColor(0xF6, 0xF7, 0xFB)     # グレーの面
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xE5, 0x53, 0x3C)     # 課題
TEAL = RGBColor(0x00, 0x94, 0x7E)      # 解決
AMBER = RGBColor(0xC2, 0x7A, 0x00)
LINE = RGBColor(0xDC, 0xDE, 0xEA)

JP = "Yu Gothic"
FOOTER_TXT = "Confidential | © 2026 Glean Technologies, Inc."

W, H = Inches(13.333), Inches(7.5)
ML = Inches(0.72)                 # 左マージン
CW = Inches(13.333 - 0.72 * 2)    # コンテンツ幅

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
_page = {"n": 0}


# ── 低レベルヘルパ ─────────────────────────────────────────────────────
def _set_typeface(run, name):
    """latin / ea / cs すべてに同じ書体を設定（日本語のフォント落ちを防ぐ）。"""
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = rPr.makeelement(qn("a:latin"), {"typeface": name})
        rPr.append(latin)
    latin.set("typeface", name)
    prev = latin
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {"typeface": name})
            prev.addnext(el)
        el.set("typeface", name)
        prev = el


def tb(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size=14, bold=False, color=BODY, space_after=4,
         space_before=0, line=1.25, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    _set_typeface(r, JP)
    return p


def rect(slide, x, y, w, h, fill, *, line_color=None, line_w=0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except (IndexError, KeyError):
            pass
    s.text_frame.word_wrap = True
    return s


def hline(slide, x, y, w, color=LINE, weight=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def arrow(slide, x, y, w, h, color=BLUE, shape=MSO_SHAPE.DOWN_ARROW):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


# ── スライド骨格 ───────────────────────────────────────────────────────
def chrome(slide, kicker, title, *, sub=None, dark=False):
    """共通ヘッダ（カテゴリ・見出し・罫線）とフッタを描く。戻り値=本文の開始 y。"""
    k_col = RGBColor(0x9A, 0xA0, 0xFF) if dark else BLUE
    t_col = WHITE if dark else INK
    s_col = RGBColor(0xB9, 0xBD, 0xD6) if dark else MUTED

    tf = tb(slide, ML, Inches(0.46), CW, Inches(0.26))
    para(tf, kicker, size=11.5, bold=True, color=k_col, space_after=0, line=1.0, first=True)

    tf = tb(slide, ML, Inches(0.76), CW, Inches(0.62))
    para(tf, title, size=26, bold=True, color=t_col, space_after=0, line=1.12, first=True)

    y = Inches(1.52)
    if sub:
        tf = tb(slide, ML, y, CW, Inches(0.3))
        para(tf, sub, size=13, color=s_col, space_after=0, line=1.25, first=True)
        y = Inches(1.9)
    hline(slide, ML, y, CW, color=RGBColor(0x3A, 0x3D, 0x66) if dark else LINE)
    footer(slide, dark=dark)
    return y + Inches(0.3)


def footer(slide, *, dark=False):
    _page["n"] += 1
    col = RGBColor(0x74, 0x78, 0x9C) if dark else RGBColor(0x9A, 0x9D, 0xB5)
    tf = tb(slide, ML, Inches(7.02), Inches(8.0), Inches(0.24))
    para(tf, FOOTER_TXT, size=8.5, color=col, space_after=0, line=1.0, first=True)
    tf = tb(slide, Inches(12.0), Inches(7.02), Inches(0.61), Inches(0.24))
    para(tf, str(_page["n"]), size=8.5, color=col, space_after=0, line=1.0,
         align=PP_ALIGN.RIGHT, first=True)


def dark_bg(slide, color=INK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def new(kicker=None, title=None, **kw):
    slide = prs.slides.add_slide(BLANK)
    y = chrome(slide, kicker, title, **kw) if kicker is not None else None
    return slide, y


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ── コンポーネント ─────────────────────────────────────────────────────
def card(slide, x, y, w, h, title, lines, *, accent=BLUE, fill=WHITE,
         label=None, title_size=14.5, body_size=11.5):
    """見出し＋箇条書きのカード。上端に細いアクセントバーを置く。"""
    rect(slide, x, y, w, h, fill, line_color=LINE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.22), y + Inches(0.24),
                                 Inches(0.055), Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False

    inner_x = x + Inches(0.38)
    inner_w = w - Inches(0.6)
    ty = y + Inches(0.19)
    if label:
        tf = tb(slide, inner_x, ty, inner_w, Inches(0.2))
        para(tf, label, size=9.5, bold=True, color=accent, space_after=0, line=1.0, first=True)
        ty += Inches(0.24)
    tf = tb(slide, inner_x, ty, inner_w, Inches(0.3))
    para(tf, title, size=title_size, bold=True, color=INK, space_after=0, line=1.15, first=True)

    body_y = ty + Inches(0.36)
    tf = tb(slide, inner_x, body_y, inner_w, h - (body_y - y) - Inches(0.2))
    for i, ln in enumerate(lines):
        para(tf, ln, size=body_size, color=BODY, space_after=5, line=1.3, first=(i == 0))


def rows(slide, y0, items, *, row_h=None, accent=BLUE, numbered=True,
         head_size=13.5, body_size=11.5, gap=0.1, bottom=6.72):
    """番号付きの横長行。items = [(見出し, 説明), ...]"""
    n = len(items)
    if row_h is None:
        row_h = (bottom - y0 / 914400 - gap * (n - 1)) / n
    y = y0
    for i, (head, detail) in enumerate(items):
        hh = Inches(row_h)
        rect(slide, ML, y, CW, hh, TINT2, line_color=None)
        if numbered:
            badge = rect(slide, ML + Inches(0.2), y + Inches(row_h / 2 - 0.16),
                         Inches(0.32), Inches(0.32), accent, shape=MSO_SHAPE.OVAL)
            tf = badge.text_frame
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            para(tf, str(i + 1), size=11, bold=True, color=WHITE, space_after=0,
                 line=1.0, align=PP_ALIGN.CENTER, first=True)
            tx = ML + Inches(0.66)
        else:
            tx = ML + Inches(0.3)
        tf = tb(slide, tx, y + Inches(0.16), Inches(2.85), hh - Inches(0.24))
        para(tf, head, size=head_size, bold=True, color=INK, space_after=0, line=1.18, first=True)
        tf = tb(slide, tx + Inches(3.0), y + Inches(0.16), CW - (tx - ML) - Inches(3.3),
                hh - Inches(0.24), anchor=MSO_ANCHOR.TOP)
        para(tf, detail, size=body_size, color=BODY, space_after=0, line=1.32, first=True)
        y += hh + Inches(gap)


def banner(slide, y, text, *, fill=TINT, color=BLUE_DK, size=12.5, h=0.52, bold=True):
    rect(slide, ML, y, CW, Inches(h), fill, line_color=None)
    tf = tb(slide, ML + Inches(0.3), y, CW - Inches(0.6), Inches(h), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, text, size=size, bold=bold, color=color, space_after=0, line=1.2, first=True)


def code_chip(slide, x, y, w, text, *, size=10, h=0.3):
    rect(slide, x, y, w, Inches(h), RGBColor(0xF2, 0xF3, 0xF9), line_color=LINE)
    tf = tb(slide, x + Inches(0.12), y, w - Inches(0.24), Inches(h), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.space_after = 0
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = BLUE_DK
    _set_typeface(r, "Consolas")


# ══════════════════════════════════════════════════════════════════════
# 1. 表紙
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.6), Inches(-1.35), Inches(4.733), Inches(4.733))
glow.fill.solid()
glow.fill.fore_color.rgb = RGBColor(0x1C, 0x1F, 0x6E)
glow.line.fill.background()
glow.shadow.inherit = False
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.05), Inches(0.075), Inches(2.5))
band.fill.solid()
band.fill.fore_color.rgb = BLUE
band.line.fill.background()
band.shadow.inherit = False

if LOGO.exists():
    s.shapes.add_picture(str(LOGO), Inches(0.9), Inches(0.75), height=Inches(0.44))

tf = tb(s, Inches(1.45), Inches(0.82), Inches(6.0), Inches(0.3))
para(tf, "Glean AI Gateway 活用デモ", size=13, bold=True, color=RGBColor(0x9A, 0xA0, 0xFF),
     space_after=0, line=1.0, first=True)

tf = tb(s, Inches(1.2), Inches(2.02), Inches(9.6), Inches(1.1))
para(tf, "AI Buddy", size=54, bold=True, color=WHITE, space_after=0, line=1.0, first=True)

tf = tb(s, Inches(1.2), Inches(3.12), Inches(9.9), Inches(1.5))
para(tf, "社内に乱立する AI エージェントを、", size=20, bold=True, color=WHITE,
     space_after=2, line=1.3, first=True)
para(tf, "社員ひとりの「唯一の入口」に束ねる", size=20, bold=True, color=WHITE,
     space_after=14, line=1.3)
para(tf, "Glean AI Gateway（LLM Gateway / MCP Gateway）＋ Personal Graph ＋ Agents API による実装デモ",
     size=13, color=RGBColor(0xB9, 0xBD, 0xD6), space_after=0, line=1.35)

for i, (lbl, x) in enumerate([("パーソナライズ", 1.2), ("エージェント・ルーティング", 3.55),
                              ("ガバナンス", 6.85), ("外部システム実行", 8.65)]):
    pill = rect(s, Inches(x), Inches(5.35), Inches([2.2, 3.15, 1.65, 2.35][i]), Inches(0.42),
                None, line_color=RGBColor(0x44, 0x48, 0x8C), radius=0.5)
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, lbl, size=11, bold=True, color=RGBColor(0xCF, 0xD2, 0xF5), space_after=0,
         line=1.0, align=PP_ALIGN.CENTER, first=True)

tf = tb(s, Inches(1.2), Inches(6.42), Inches(10.0), Inches(0.26))
para(tf, "Solutions Engineering / 2026", size=11, color=RGBColor(0x74, 0x78, 0x9C),
     space_after=0, line=1.0, first=True)
footer(s, dark=True)
notes(s, "AI Buddy は、部署ごとにバラバラの基盤で作られた社内エージェントを、社員から見て"
         "ひとつの入口に束ねるデモ。モデル呼び出しとツール実行はすべて Glean AI Gateway を通る。")

# ══════════════════════════════════════════════════════════════════════
# 2. エグゼクティブサマリ
# ══════════════════════════════════════════════════════════════════════
s, y = new("エグゼクティブサマリ", "デモで示すのは「探す・聞く・やる」が1画面で完結する体験",
           sub="社員は AI Buddy に話しかけるだけ。裏側の基盤の違い・権限・モデル選定はプラットフォームが吸収する。")
cw = Emu(int((CW - Inches(0.3) * 2) / 3))
for i, (label, title, lines, accent) in enumerate([
    ("① パーソナライズ", "あなたの文脈から始まる",
     ["Glean Personal Graph（memory ツール）を起動時にライブ取得。",
      "進行中プロジェクト・未処理スレッドを踏まえた先回り提案を提示。",
      "「さっきの PR」だけで対象を特定できる。"], BLUE),
    ("② ルーティング", "23 体を横断して最適な1体へ",
     ["Copilot Studio / Dify / ローコード / GitHub Actions / 内製 が混在。",
      "所属部署・permission・評価★・実行成功率で選定。",
      "選定理由を思考トレースとして可視化。"], TEAL),
    ("③ 実行", "相談で終わらせず、書き込みまで",
     ["MCP Gateway 経由で Jira 起票・Slack 投稿・メール送信・GitHub 操作。",
      "「法務に確認 → Jira 起票 → Slack 通知」を1依頼で完遂。",
      "ツールはソースの権限を継承。"], AMBER),
]):
    card(s, ML + Emu(int(i * (cw + Inches(0.3)))), y, cw, Inches(3.55), title, lines,
         accent=accent, label=label)
banner(s, Inches(6.05),
       "すべてのモデル呼び出しとツール実行が Glean AI Gateway を通る ＝ 統制・可視化・監査が1か所に集まる",
       h=0.6)
notes(s, "3枚のカードがそのままデモの流れ。最後の帯が AI Gateway の価値（コントロールプレーン）。")

# ══════════════════════════════════════════════════════════════════════
# 3. 課題（マクロ）
# ══════════════════════════════════════════════════════════════════════
s, y = new("顧客の課題 ①", "AI スプロールは、すでに企業の「新常態」",
           sub="AI ツールは組織全体へ広がる一方で、共通のガバナンス層が存在しない。")
cw2 = Emu(int((CW - Inches(0.28)) / 2))
ch = Inches(2.02)
for i, (label, title, lines) in enumerate([
    ("AI SPRAWL", "AI スプロール",
     ["Claude / Cursor / Copilot、各部署の自作エージェントが並列で増殖。",
      "IT が把握する前に現場が使い始め、棚卸しができない。"]),
    ("COST & CONTROLS", "コストと統制",
     ["誰がどのモデルにいくら使っているか、部署・アプリ単位で見えない。",
      "ルーティングも予算上限もなく、請求が来るまで気づけない。"]),
    ("SECURITY & GOVERNANCE", "セキュリティとガバナンス",
     ["AI の面ごとに独自ルール。ポリシーより先に新ツールが到着する。",
      "どのツールがどの引数で呼ばれたか、実行時がブラックボックス。"]),
    ("KNOWLEDGE FRAGMENTATION", "知識の分断",
     ["各 AI アプリがゼロから始まり、共有メモリを持たない。",
      "ツールを増やすほど文脈が散り、賢さが積み上がらない。"]),
]):
    x = ML + Emu(int((i % 2) * (cw2 + Inches(0.28))))
    yy = y + Emu(int((i // 2) * (ch + Inches(0.24))))
    card(s, x, yy, cw2, ch, title, lines, accent=CORAL, label=label)
banner(s, Inches(6.1),
       "出典: Glean「AI Gateway FAQ」— AI スプロールが、コスト管理・セキュリティポリシー・知識共有の3点で穴を生む",
       fill=TINT2, color=MUTED, size=10.5, h=0.42, bold=False)
notes(s, "Glean 公式の AI Gateway ポジショニングに沿った4つの課題。"
         "顧客の現状ヒアリングをここに当てはめる。")

# ══════════════════════════════════════════════════════════════════════
# 4. 課題（現場）
# ══════════════════════════════════════════════════════════════════════
s, y = new("顧客の課題 ②", "社内エージェントを増やすほど、社員は使わなくなる",
           sub="エージェントは作られている。問題は「見つからない・選べない・実行まで届かない」。")
rows(s, y, [
    ("入口が分からない",
     "法務は Copilot Studio、研究は Dify、営業は内製、CI/CD は GitHub Actions。基盤ごとに UI も URL も別で、社員は「どこに聞けばいいか」を知らない。"),
    ("権限がバラバラ",
     "権限モデルが基盤ごとに異なり、部署をまたぐ相談が可能かどうか本人には判断できない。結果、越境相談そのものが起きなくなる。"),
    ("品質が見えない",
     "同じ仕事ができるエージェントが複数あっても、どれが実績があるのか分からない。新しい方が良いのか、古い方が安全なのかの判断材料がない。"),
    ("文脈がゼロから",
     "毎回「自分は誰で、いま何のプロジェクトをしていて、どの PR の話か」を説明し直す。曖昧な指示は通らないので、質問を書くコスト自体が高い。"),
    ("相談で終わる",
     "回答は得られるが、Jira 起票・Slack 連絡・メール送信は結局手作業。「AI に聞いた後」の作業時間が減らない。"),
], accent=CORAL, gap=0.11)
notes(s, "①はマクロ（IT/経営の言葉）、②は現場の言葉。顧客の反応が良い方を深掘りする。")

# ══════════════════════════════════════════════════════════════════════
# 5. セクション：ソリューション
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.6), Inches(0.075), Inches(1.9))
band.fill.solid()
band.fill.fore_color.rgb = BLUE
band.line.fill.background()
band.shadow.inherit = False
tf = tb(s, Inches(1.2), Inches(2.35), Inches(4.0), Inches(0.3))
para(tf, "SOLUTION", size=12, bold=True, color=RGBColor(0x9A, 0xA0, 0xFF), space_after=0,
     line=1.0, first=True)
tf = tb(s, Inches(1.2), Inches(2.78), Inches(10.6), Inches(1.9))
para(tf, "AI Buddy", size=40, bold=True, color=WHITE, space_after=10, line=1.05, first=True)
para(tf, "社員から見た入口はひとつ。基盤の違い・権限・モデル選定という複雑さは、"
         "Glean AI Gateway が吸収する。", size=17, color=RGBColor(0xC7, 0xCA, 0xE8),
     space_after=0, line=1.4)
footer(s, dark=True)
notes(s, "ここから解決策。既存のエージェント資産を捨てずに束ねる、という点を強調。")

# ══════════════════════════════════════════════════════════════════════
# 6. AI Buddy の4層
# ══════════════════════════════════════════════════════════════════════
s, y = new("ソリューション", "AI Buddy は「文脈 → 選定 → 統制 → 実行」の4層で動く",
           sub="1回の依頼の中で4層すべてが動き、その過程はチャット上のトレースとして開示される。")
rows(s, y, [
    ("① パーソナライズ",
     "Glean Personal Graph（MCP の memory ツール）から進行中プロジェクト・未処理スレッド・主要コラボレーターをライブ取得し、起動時に先回り提案を提示する。"),
    ("② ルーティング",
     "質問意図・所属部署・レジストリのメタ情報から、最適な専門エージェントを選定。Glean LLM Gateway 上の tool-calling として実装。"),
    ("③ ガバナンス",
     "permission（利用可能部署）の最終判定をバックエンドで実施。さらに評価★と実行成功率で「同じ仕事ができる複数候補」から実績のある1体を選ぶ。"),
    ("④ 実行",
     "MCP Gateway のツールで Jira 起票・Slack 投稿・メール送信・GitHub ブランチ / PR 作成まで完遂。相談から実行までを1依頼で閉じる。"),
], accent=BLUE, gap=0.13)
notes(s, "デモシナリオ集の ①〜⑤ に対応。⑤は①〜④の合わせ技。")

# ══════════════════════════════════════════════════════════════════════
# 7. アーキテクチャ
# ══════════════════════════════════════════════════════════════════════
s, y = new("アーキテクチャ", "自作 UI から見た Glean は「4つの API 面」だけ",
           sub="アプリ側にモデルの API キーもツールの認証情報も持たない。すべて Glean 経由。")

# 上段: UI
rect(s, ML, Inches(2.02), Inches(4.3), Inches(0.6), INK, line_color=None)
tf = tb(s, ML + Inches(0.25), Inches(2.02), Inches(3.8), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "AI Buddy UI（React + Vite）", size=12.5, bold=True, color=WHITE, space_after=0,
     line=1.0, first=True)
arrow(s, ML + Inches(1.95), Inches(2.7), Inches(0.4), Inches(0.34), color=RGBColor(0xC2, 0xC5, 0xE0))
rect(s, ML, Inches(3.12), Inches(4.3), Inches(0.6), TINT, line_color=BLUE)
tf = tb(s, ML + Inches(0.25), Inches(3.12), Inches(3.8), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "オーケストレータ（Express / buddy.js）", size=12.5, bold=True, color=BLUE_DK,
     space_after=0, line=1.0, first=True)

# OAuth
rect(s, ML, Inches(4.02), Inches(4.3), Inches(0.94), TINT2, line_color=LINE)
tf = tb(s, ML + Inches(0.25), Inches(4.16), Inches(3.8), Inches(0.7))
para(tf, "OAuth 2.0（Authorization Code + PKCE）", size=11.5, bold=True, color=INK,
     space_after=3, line=1.15, first=True)
para(tf, "SEARCH / CHAT / AGENTS / DOCUMENTS /\nENTITIES / TOOLS / MCP / LLM_PROXY",
     size=9.5, color=MUTED, space_after=0, line=1.25)

# 右: Glean AI Gateway コンテナ
gx, gy, gw, gh = Inches(5.45), Inches(2.02), Inches(7.16), Inches(2.94)
rect(s, gx, gy, gw, gh, TINT, line_color=BLUE, line_w=1.25)
tf = tb(s, gx + Inches(0.28), gy + Inches(0.2), gw - Inches(0.56), Inches(0.26))
para(tf, "Glean AI Gateway ＋ プラットフォーム API", size=12.5, bold=True, color=BLUE_DK,
     space_after=0, line=1.0, first=True)
sub_w = Emu(int((gw - Inches(0.56) - Inches(0.2)) / 2))
for i, (name, detail, col) in enumerate([
    ("LLM Gateway", "モデル層。Messages API /\nResponses API に Bearer 転送", BLUE),
    ("MCP Gateway", "アクション層。tools/list で発見、\ntools/call で実行", TEAL),
    ("Personal Graph", "memory ツールで個人の\n業務文脈をライブ取得", AMBER),
    ("Agents API", "agents/search と runs/stream で\n実在エージェントを実行", RGBColor(0x7A, 0x3D, 0xC4)),
]):
    x = gx + Inches(0.28) + Emu(int((i % 2) * (sub_w + Inches(0.2))))
    yy = gy + Inches(0.56) + Emu(int((i // 2) * (Inches(1.06) + Inches(0.14))))
    rect(s, x, yy, sub_w, Inches(1.06), WHITE, line_color=LINE)
    tf = tb(s, x + Inches(0.18), yy + Inches(0.15), sub_w - Inches(0.36), Inches(0.8))
    para(tf, name, size=12, bold=True, color=col, space_after=3, line=1.0, first=True)
    para(tf, detail, size=9.5, color=BODY, space_after=0, line=1.22)
arrow(s, Inches(4.98), Inches(3.24), Inches(0.4), Inches(0.34), color=BLUE,
      shape=MSO_SHAPE.RIGHT_ARROW)

# 下段: A2A → 部署別エージェント
arrow(s, ML + Inches(1.95), Inches(5.06), Inches(0.4), Inches(0.34), color=RGBColor(0xC2, 0xC5, 0xE0))
rect(s, ML, Inches(5.48), CW, Inches(1.1), TINT2, line_color=LINE)
tf = tb(s, ML + Inches(0.3), Inches(5.62), CW - Inches(0.6), Inches(0.85))
para(tf, "A2A（message/send・agent-card.json）で部署別エージェントを統一呼び出し", size=12,
     bold=True, color=INK, space_after=4, line=1.1, first=True)
para(tf, "法務 / 人事 / 研究 / 営業 / IT / エンジニアリング / QA / SRE / データ / セキュリティ / 財務 / 経営　"
         "— 基盤は Copilot Studio・Dify・ローコード・GitHub Actions・Jenkins・内製が混在（計 23 体）",
     size=10.5, color=BODY, space_after=0, line=1.3)
notes(s, "重要点は2つ。(1) アプリはプロバイダのキーを持たない。(2) 既存の他社基盤は A2A で"
         "そのまま繋がる（作り直し不要）。")

# ══════════════════════════════════════════════════════════════════════
# 8. Glean AI Gateway とは
# ══════════════════════════════════════════════════════════════════════
s, y = new("Glean AI Gateway", "AI Gateway は LLM Gateway と MCP Gateway で構成される",
           sub="「AI が何にアクセスできるか」と「AI が何をできるか」を、それぞれ1つの統制層に集約する。")
cw2 = Emu(int((CW - Inches(0.3)) / 2))
card(s, ML, y, cw2, Inches(2.6), "LLM Gateway — モデル層",
     ["スコープ: モデルへのアクセスと推論。",
      "どのモデルを使えるか、予算をどう配分するかを制御。",
      "30+ のフロンティア / OSS モデルに1つの統制層から到達。",
      "クライアント（Claude Code, Codex 等）と LLM プロバイダの間に立つ透過プロキシ。"],
     accent=BLUE, label="MODEL LAYER")
card(s, ML + Emu(int(cw2 + Inches(0.3))), y, cw2, Inches(2.6), "MCP Gateway — アクション層",
     ["スコープ: ツール実行・データソース・アクション。",
      "エージェントがどのツールに触れられるかを制御。",
      "2,000+ の統制済みツールを1エンドポイントで公開。",
      "ツールはソースシステムの権限を継承し、Protect+ のガードレールが前段に立つ。"],
     accent=TEAL, label="ACTION LAYER")
rows(s, Inches(4.98), [
    ("可視化", "モデル別・ユーザー別・アプリ別にリクエスト量とトークン消費を分解。どの MCP サーバのどのツールが誰に使われているかが見える。"),
    ("統制", "ユーザー / アプリ単位のモデル・クォータ、低コストモデルへのフォールバック、全モデル・全ツールへのアクセス制御。"),
], accent=BLUE_DK, numbered=False, row_h=0.82, gap=0.1)
notes(s, "この2層の分担が AI Gateway の骨格。デモでは LLM GW と MCP GW の両方を実際に叩いている。")

# ══════════════════════════════════════════════════════════════════════
# 9. 組み込み① LLM Gateway
# ══════════════════════════════════════════════════════════════════════
s, y = new("組み込み ①  LLM Gateway", "AI Buddy の「司令塔モデル」は 100% Glean 経由",
           sub="ルーティング判断そのものが Gateway 上の tool-calling。モデルを差し替えてもアプリ改修はゼロ。")
rows(s, y, [
    ("エンドポイント",
     "Claude 系は POST /rest/api/v1/anthropic/v1/messages（Messages API）、GPT 系は POST /rest/api/v1/openai/v1/responses（Responses API）。業界標準仕様に準拠しているため実装が使い回せる。"),
    ("認証",
     "Glean OAuth トークンを Bearer 注入（scope: LLM_PROXY）。Anthropic / OpenAI の API キーはアプリ側に一切持たない。トークンはサーバ側セッションにのみ保持しレスポンスへ含めない。"),
    ("モデル選択",
     "GET /rest/api/v1/{openai|anthropic|gemini}/v1/models でカタログをライブ取得。UI の「テスト＆有効化」で1回実呼び出しして疎通確認したモデルだけを司令塔候補に追加できる。"),
    ("フォールバック",
     "テナントの提供モデルは変動するため、モデル起因のエラー時のみ候補を順に試し、成功したモデルをキャッシュ。プロバイダ障害・モデル廃止でデモが止まらない。"),
    ("実装ファイル",
     "backend/src/routes/llmGateway.js（Gateway プロキシ・モデル管理）、backend/src/lib/buddy.js（tool-calling でのルーティングと要約）。"),
], accent=BLUE, gap=0.1)
notes(s, "「ロックイン回避」と「コストの可視化」がここで効く。顧客が既に Claude Code や Codex を"
         "使っているなら、同じ Gateway にそれも寄せられる、と繋げる。")

# ══════════════════════════════════════════════════════════════════════
# 10. 組み込み② MCP Gateway
# ══════════════════════════════════════════════════════════════════════
s, y = new("組み込み ②  MCP Gateway", "「相談」で終わらせず、外部システムへの実行まで届かせる",
           sub="ツールを固定実装せず、tools/list で動的に発見して LLM に提示する。増えたツールは自動で使えるようになる。")
rows(s, y, [
    ("エンドポイント",
     "{instanceUrl}/mcp/{server}（Streamable HTTP + JSON-RPC）。initialize → tools/list → tools/call の3ステップのみ。ツールごとの個別実装は不要。"),
    ("実行できること",
     "Jira 課題作成 / Slack メッセージ投稿 / Gmail 送信 / GitHub ブランチ・PR 作成 / エンタープライズ検索 / ドキュメント取得。Glean 側でツールを追加すればアプリを触らずに増える。"),
    ("権限とセキュリティ",
     "ツールはソースシステムの権限を継承し、ユーザーが元々持つ以上の権限は出ない。Protect+ のプロンプトインジェクション検知・悪性コード検知・整合性チェックが全ツール呼び出しの前段に立つ。"),
    ("Buddy 側の作り込み",
     "発見したツールをそのまま LLM の tool として提示。「送って / 起票して」と明示されたら確認で止まらず実行し、不足パラメータは personal memory と直前のエージェント出力から補完する。"),
    ("実装ファイル",
     "backend/src/lib/mcpClient.js（MCP クライアント）、frontend/src/components/ToolsView.jsx（ツール一覧・単体テスト実行）。"),
], accent=TEAL, gap=0.1)
notes(s, "デモでは実テナントに反映される（Jira 起票・Slack 投稿）。事前にテスト用の"
         "プロジェクト / チャンネルを確認しておく。")

# ══════════════════════════════════════════════════════════════════════
# 11. 組み込み③ Personal Graph
# ══════════════════════════════════════════════════════════════════════
s, y = new("組み込み ③  Personal Graph（memory）", "「あなただから、この曖昧な指示が通る」",
           sub="合成データは使わない。接続したユーザー本人の Glean Personal Graph をライブで読む。")
cw2 = Emu(int((CW - Inches(0.3)) / 2))
card(s, ML, y, cw2, Inches(2.64), "取得しているもの",
     ["進行中プロジェクト / 未処理スレッド",
      "主要コラボレーター / 最近の話題",
      "仕事の進め方の好み（preference signals）",
      "→ system プロンプトへ contextBlock として注入",
      "実装: backend/src/lib/memory.js"],
     accent=AMBER, label="MCP memory ツール")
card(s, ML + Emu(int(cw2 + Inches(0.3))), y, cw2, Inches(2.64), "デモで起きること",
     ["「さっき作ったドラフト PR、レビュー観点を整理して」",
      "→ memory から PR #80 / #81 / #82 を特定",
      "→ CI/CD・GitHub エージェントへ自動ルーティング",
      "起動時には未処理スレッド由来の先回り提案カードを生成",
      "主語・目的語が無い依頼でも成立する"],
     accent=AMBER, label="曖昧な指示が通る")
banner(s, Inches(5.06),
       "共有メモリが「知識の分断」を解く — Glean が持たない AI の面にも、同じ enterprise context を配れる",
       h=0.56)
rows(s, Inches(5.84), [
    ("前提条件", "テナントの MCP サーバが memory ツールを公開している必要がある（未公開の場合は UI にエラーを表示し、合成データでの代替はしない）。"),
], accent=MUTED, numbered=False, row_h=0.8, gap=0.0)
notes(s, "パーソナライズは「聞かれたことに答える」から「先に出す」への転換。"
         "Proactive Intelligence の文脈に接続できる。")

# ══════════════════════════════════════════════════════════════════════
# 12. 組み込み④ Agents API + A2A
# ══════════════════════════════════════════════════════════════════════
s, y = new("組み込み ④  Agents API ＋ A2A", "マルチ基盤に散らばるエージェントを、ひとつのカタログへ",
           sub="既存の Copilot Studio / Dify 資産を作り直さない。A2A ラッパーを立てて統合カタログに載せる。")
rows(s, y, [
    ("統合レジストリ",
     "23 体を description / permission / metadata / 評価★ / 実行成功率 付きの統合カタログとして表示。社員は基盤の違いを意識せず、Buddy が代わりに選ぶ。"),
    ("基盤の多様性を吸収",
     "Copilot Studio・Dify・ローコード・GitHub Actions・Jenkins・内製が混在。A2A（/.well-known/agent-card.json + message/send）で呼び出し方を統一する。"),
    ("実在エージェントも混在",
     "POST /rest/api/v1/agents/search でテナントの実在 Glean エージェントを検索してレジストリへ追加。ルーティングされると runs/stream で本物が実行され、トレースに「Glean（実物）」バッジが付く。"),
    ("信頼性シグナルで選ぶ",
     "同機能が3つ並ぶとき（★4.8 / 成功率95%、★3.8 / 63%、★4.9 / 58%）、Buddy は実行成功率を優先。「★は高いが実績が浅い Beta を避ける」判断を思考として明示する。"),
    ("実装ファイル",
     "examples/agents/domainSpecs.mjs（エージェント定義）、backend/src/lib/a2a.js（A2A クライアント）、backend/src/lib/gleanAgents.js（実在エージェント検索・実行）。"),
], accent=RGBColor(0x7A, 0x3D, 0xC4), gap=0.1)
notes(s, "「評価≠信頼性」を体感させるのがこのスライドの狙い。"
         "顧客が既に複数基盤を持っている場合、最も刺さるパート。")

# ══════════════════════════════════════════════════════════════════════
# 13. ガバナンス
# ══════════════════════════════════════════════════════════════════════
s, y = new("ガバナンス", "1回の OAuth で、モデル・ツール・エージェント・検索まで通る",
           sub="認証・認可・監査の起点が1つに集まるため、AI の面が増えても統制コストが増えない。")
rows(s, y, [
    ("認証方式",
     "OAuth 2.0 Authorization Code + PKCE（Glean 認可サーバ）。/.well-known/oauth-authorization-server から自動解決。トークンはサーバ側セッションのみ、レスポンスには一切含めない。"),
    ("スコープ",
     "SEARCH / CHAT / AGENTS / DOCUMENTS / ENTITIES / TOOLS / MCP / LLM_PROXY。この1セットでモデル呼び出し・ツール実行・エージェント実行・検索がすべてカバーされる。"),
    ("権限判定の置き場所",
     "permission（利用可能部署）の最終判定はバックエンドで行い、LLM の判断には委ねない。権限が無い場合は理由と正しい依頼先を利用者に説明する。"),
    ("観測",
     "全トラフィックが Gateway を通るため、部署別・アプリ別・ユーザー別の消費と、どの MCP ツールが誰に呼ばれたかを Insights で追跡できる。"),
], accent=BLUE_DK, gap=0.12)
banner(s, Inches(6.0),
       "デモでの確認例: 所属を「IT」にして研究系を質問 → 研究エージェントは権限なしと判定され、Buddy が代替の依頼先を案内する",
       fill=TINT2, color=BODY, size=11, h=0.5, bold=False)
notes(s, "権限判定をアプリ側（バックエンド）に置いている点は必ず触れる。"
         "LLM に権限を判断させない設計はセキュリティ担当に効く。")

# ══════════════════════════════════════════════════════════════════════
# 14. デモシナリオ
# ══════════════════════════════════════════════════════════════════════
s, y = new("デモの見せ方", "5段階で「効く要素」を積み上げる",
           sub="下の段ほど判断材料が増える。顧客の関心に応じて ③ か ⑤ から入るのも有効。")
rows(s, y, [
    ("所属部署だけで変わる",
     "「経費精算の上限と承認フローを教えて」を営業／エンジニアリングで送信 → 全社ポリシー＋役職別を統合した別々の回答になる。"),
    ("＋ Personal Memory",
     "「さっき作ったドラフト PR、レビュー観点を整理して」→ memory が PR #80-82 を補完し、指示が曖昧でも成立する。"),
    ("＋ 信頼性シグナル",
     "「一番確実な方法で本番リリースして」→「確実＝実行成功率」と解釈し、★は低くても成功率95%のエージェントを選ぶ。"),
    ("＋ Tools 実行",
     "「今日の PR レビュー状況をチームの Slack にまとめて投稿して」→ MCP Gateway 経由で実テナントに反映される。"),
    ("一気通貫（本命）",
     "「NDA でひな形から外れた条項を法務に確認 → 修正点を Jira に起票 → 営業担当に Slack 通知」を1依頼で完遂する。"),
], accent=TEAL, gap=0.1)
notes(s, "詳細プロンプトは DEMO_SCENARIOS.md に全掲載。時間が10分なら ② → ⑤ の2本に絞る。")

# ══════════════════════════════════════════════════════════════════════
# 15. 価値と次のステップ
# ══════════════════════════════════════════════════════════════════════
s, y = new("顧客価値", "3つの立場それぞれに、別々の理由で効く",
           sub=None)
cw = Emu(int((CW - Inches(0.3) * 2) / 3))
for i, (label, title, lines, accent) in enumerate([
    ("現場の社員", "探す時間がゼロになる",
     ["入口は1つ。どの基盤に何があるかを覚えなくていい。",
      "曖昧な指示でも文脈が補完される。",
      "相談から Jira / Slack / メールまで1依頼で終わる。"], BLUE),
    ("IT / プラットフォーム", "既存資産を捨てずに統合",
     ["Copilot Studio・Dify・内製をそのまま A2A で載せる。",
      "権限・監査・ガードレールが1か所に集まる。",
      "ツールを Glean 側で足せば、アプリ改修なしで増える。"], TEAL),
    ("経営 / 財務", "消費が見える・ロックインしない",
     ["モデル・ツールの消費を部署別／アプリ別に可視化。",
      "30+ モデルから選べるためベンダー依存を避けられる。",
      "低コストモデルへのフォールバックで単価を下げられる。"], AMBER),
]):
    card(s, ML + Emu(int(i * (cw + Inches(0.3)))), Inches(1.9), cw, Inches(2.78),
         title, lines, accent=accent, label=label)

tf = tb(s, ML, Inches(4.84), CW, Inches(0.3))
para(tf, "次のステップ（提案）", size=14, bold=True, color=INK, space_after=0, line=1.0, first=True)
rows(s, Inches(5.24), [
    ("PoC スコープ", "貴社の既存エージェント2〜3体を A2A / MCP でこのレジストリに接続し、社員から見た入口を1つにする。"),
    ("パーソナライズ", "Personal Graph を使った先回り提案の対象業務を1つ選定する（例: PR レビュー、案件フォロー）。"),
    ("Gateway 移行", "LLM Gateway 経由に寄せる AI の面を決める（Claude Code / Codex / 自社アプリ）。"),
], accent=BLUE, numbered=False, row_h=0.5, gap=0.08)
notes(s, "顧客の立場に応じて読み上げるカードを変える。次のステップは合意を取るための叩き台。")

# ══════════════════════════════════════════════════════════════════════
# 16. 補足・前提
# ══════════════════════════════════════════════════════════════════════
s, y = new("補足", "デモの前提と、実装の割り切りポイント",
           sub="顧客に誤解を与えないため、合成部分とライブ部分を明示する。")
rows(s, y, [
    ("評価★／実行成功率",
     "Glean 標準機能にはない概念で、本デモ用の合成メタ（examples/agents/domainSpecs.mjs）。UI 上でも「合成メタ」と明示している。実運用では Agent Governance の指標に置き換える。"),
    ("部署別エージェント",
     "A2A 準拠のモックサーバ（ポート 5601-5625）。実運用では Copilot Studio / Dify を A2A 準拠に寄せるか、A2A ラッパーを立てる想定。モック側は一切「Glean」を名乗らない。"),
    ("Personal Memory",
     "接続ユーザーの Glean Personal Graph をライブ取得（合成データなし）。テナントの MCP サーバが memory ツールを公開している必要がある。"),
    ("Tools の実行",
     "Jira / Slack / Gmail / GitHub は実テナントに反映される。灰色の DEMO バッジ付きツール（Salesforce 等）は外部接続せず擬似応答を返す。"),
    ("テナント前提",
     "LLM Gateway が有効（LlmFabric_Enabled）で、OAuth スコープに LLM_PROXY と MCP を含めて接続済みであること。"),
], accent=MUTED, gap=0.1)
notes(s, "SE として必ず伝える前提。特に「★は合成」「モックは Glean を名乗らない」は誤解防止のため必須。")

# ══════════════════════════════════════════════════════════════════════
out = Path(sys.argv[1] if len(sys.argv) > 1 else ROOT / "AI_Buddy_デモ紹介.pptx")
prs.save(str(out))
print(f"saved: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
