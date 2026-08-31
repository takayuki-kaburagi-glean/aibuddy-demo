#!/usr/bin/env python3
"""pptx のテキストはみ出しを機械的に検査する（CJK 幅を考慮した概算）。

usage: PYTHONPATH=... python3 scripts/qa_deck.py deck.pptx
"""
import sys
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

EMU_IN = 914400


def char_w(ch):
    """1文字の幅（em 単位の概算）。全角=1.0、半角英数=0.52、空白=0.28。"""
    if ch in " \u3000":
        return 1.0 if ch == "\u3000" else 0.28
    return 1.0 if unicodedata.east_asian_width(ch) in ("W", "F", "A") else 0.52


def est_lines(text, size_pt, box_w_in):
    """折り返し後の行数を概算（明示改行も考慮）。"""
    if not text:
        return 1
    usable_em = (box_w_in * 72.0) / size_pt
    total = 0
    for seg in text.split("\n"):
        w = sum(char_w(c) for c in seg)
        total += max(1, -(-int(w * 100) // int(usable_em * 100)) if usable_em > 0 else 1)
    return total


def scan(path):
    prs = Presentation(str(path))
    problems = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            if not tf.text.strip():
                continue
            bw = (sh.width - (tf.margin_left or 0) - (tf.margin_right or 0)) / EMU_IN
            bh = (sh.height - (tf.margin_top or 0) - (tf.margin_bottom or 0)) / EMU_IN
            if bw <= 0:
                continue
            need = 0.0
            for p in tf.paragraphs:
                txt = "".join(r.text for r in p.runs)
                size = None
                for r in p.runs:
                    if r.font.size:
                        size = r.font.size.pt
                        break
                size = size or 18.0
                ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.2
                sa = p.space_after.pt if p.space_after else 0
                sb = p.space_before.pt if p.space_before else 0
                need += (est_lines(txt, size, bw) * size * ls * 1.22 + sa + sb) / 72.0
            # 右端まで使い切っている行は 1 行分の余白を許容
            if need > bh + 0.02:
                problems.append((i, round(need, 2), round(bh, 2),
                                 tf.text.replace("\n", " / ")[:80]))
    return prs, problems


def main():
    path = Path(sys.argv[1])
    prs, problems = scan(path)
    print(f"{path.name}: {len(prs.slides._sldIdLst)} slides, "
          f"{prs.slide_width / EMU_IN:.2f}x{prs.slide_height / EMU_IN:.2f} in")
    # スライドの下端(7.5in)を超えて配置された図形も検出
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            bottom = (sh.top + sh.height) / EMU_IN
            right = (sh.left + sh.width) / EMU_IN
            if bottom > 7.52 or right > 13.36:
                label = (sh.text_frame.text[:40] if sh.has_text_frame else sh.shape_type)
                print(f"  [OUT OF BOUNDS] s{i}: bottom={bottom:.2f} right={right:.2f} {label!r}")
    if not problems:
        print("  overflow: none detected")
    for i, need, have, txt in problems:
        print(f"  [OVERFLOW] s{i}: need {need}in > box {have}in :: {txt!r}")
    print(f"  total overflow candidates: {len(problems)}")


if __name__ == "__main__":
    main()
